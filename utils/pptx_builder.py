from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml import parse_xml
import pandas as pd
from typing import Dict, List, Tuple
import re

# ==========================================
# 🎨 ФИРМЕННЫЕ ЦВЕТА (из Брендбука)
# ==========================================
class PJColors:
    TOMATO_RED = RGBColor(225, 45, 38)
    BASIL_GREEN = RGBColor(3, 89, 45)
    OLIVE_GREEN = RGBColor(17, 72, 52)
    CHEESE_YELLOW = RGBColor(255, 225, 33)
    DOUGH_BEIGE = RGBColor(244, 232, 220)
    BAKED_RED = RGBColor(127, 22, 31)
    WHITE = RGBColor(255, 255, 255)
    BLACK = RGBColor(0, 0, 0)

class PJFonts:
    HEADLINE = "Oswald"
    BODY = "Roboto Condensed"

class PJSlogans:
    MAIN = "Лучшие ингредиенты. Лучшая пицца."
    EMOTION = "Когда есть вкус, есть эмоции."
    TRANSPARENCY = "Мы открыты, честны и уверены в своем продукте."
    CARE = "Дарить людям хорошее настроение через еду и сервис."

# ==========================================
# 🛠️ ВСПОМОГАТЕЛЬНЫЕ ФУНКЦИИ
# ==========================================
def _set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def _add_text(slide, left, top, width, height, text, font_name, font_size, color, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def _format_cell(cell, text, font_name, font_size, font_color, bg_color=None, bold=False, align=PP_ALIGN.CENTER):
    if pd.isna(text) or text is None:
        text = ""
    cell.text = str(text)
    
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.name = font_name
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = font_color
        paragraph.font.bold = bold
        paragraph.alignment = align
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color

    tcPr = cell._tc.get_or_add_tcPr()
    for border in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
        ln = parse_xml(f'<{border} w="0" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><a:noFill/></{border}>')
        tcPr.append(ln)

# ==========================================
# 🧠 УМНЫЙ ПАРСЕР EXCEL ФАЙЛОВ КР
# ==========================================
def parse_kr_excel(file_path) -> Dict[str, pd.DataFrame]:
    """
    Парсит Excel файл Клиентского Рейтинга и извлекает структурированные данные.
    """
    # Читаем все листы
    all_sheets = pd.read_excel(file_path, sheet_name=None, header=None)
    
    result = {}
    
    for sheet_name, df in all_sheets.items():
        # Заполняем NaN
        df = df.fillna('')
        
        if 'Оценки' in str(sheet_name) or 'Оценки' in str(df.iloc[0, 0]):
            # Парсим лист с оценками
            result['сайт_спб'] = _extract_ratings_block(df, 'Сайт', 'СПБ')
            result['сайт_тюмень'] = _extract_ratings_block(df, 'Сайт', 'Тюмень')
            result['агрегаторы_спб'] = _extract_ratings_block(df, 'Агрегаторы', 'СПБ')
            result['агрегаторы_тюмень'] = _extract_ratings_block(df, 'Агрегаторы', 'Тюмень')
            result['геосервисы_спб'] = _extract_ratings_block(df, 'Геосервисы', 'СПБ')
            result['геосервисы_тюмень'] = _extract_ratings_block(df, 'Геосервисы', 'Тюмень')
            
        elif 'Анализ отзывов' in str(sheet_name) or 'Анализ отзывов' in str(df.iloc[0, 0]):
            # Парсим лист с анализом отзывов
            result['жалобы'] = _extract_reviews_block(df, 'жалобы')
            result['позитив'] = _extract_reviews_block(df, 'позитив')
    
    return result

def _extract_ratings_block(df, source_type: str, city: str) -> pd.DataFrame:
    """
    Извлекает блок оценок для конкретного источника и города.
    """
    # Находим начало блока
    start_row = None
    for idx, row in df.iterrows():
        cell_value = str(row[0]).strip()
        if source_type in cell_value and city in str(row.iloc[1:]).strip():
            # Проверяем, что это заголовок блока (Сайт/Агрегаторы/Геосервисы + город)
            if any(keyword in cell_value for keyword in ['Сайт', 'Агрегаторы', 'Геосервисы']):
                start_row = idx
                break
    
    if start_row is None:
        return pd.DataFrame()
    
    # Ищем строки с данными ресторанов
    data_rows = []
    for idx in range(start_row + 1, len(df)):
        row = df.iloc[idx]
        restaurant = str(row[0]).strip()
        
        # Пропускаем пустые строки и итоги
        if not restaurant or 'Итого' in restaurant or 'Кол-во' in restaurant:
            continue
        
        # Проверяем, что это название ресторана (не число)
        if restaurant and not restaurant.replace('.', '').replace(',', '').isdigit():
            # Извлекаем данные: колонки 1-5 (оценки), всего, средний рейтинг
            try:
                rating_1 = row[1] if pd.notna(row[1]) else 0
                rating_2 = row[2] if pd.notna(row[2]) else 0
                rating_3 = row[3] if pd.notna(row[3]) else 0
                rating_4 = row[4] if pd.notna(row[4]) else 0
                rating_5 = row[5] if pd.notna(row[5]) else 0
                total = row[6] if pd.notna(row[6]) else ''
                avg_rating = row[7] if pd.notna(row[7]) else ''
                
                data_rows.append({
                    'Ресторан': restaurant,
                    '1': rating_1,
                    '2': rating_2,
                    '3': rating_3,
                    '4': rating_4,
                    '5': rating_5,
                    'Всего': total,
                    'Средний рейтинг': avg_rating
                })
            except:
                continue
    
    return pd.DataFrame(data_rows)

def _extract_reviews_block(df, review_type: str) -> pd.DataFrame:
    """
    Извлекает блок анализа отзывов (жалобы или позитив).
    """
    # Находим начало блока
    start_row = None
    for idx, row in df.iterrows():
        cell_value = str(row[0]).strip()
        if 'Анализ отзывов' in cell_value and review_type in cell_value:
            start_row = idx
            break
    
    if start_row is None:
        return pd.DataFrame()
    
    # Определяем заголовки колонок
    if review_type == 'жалобы':
        headers = ['Ресторан', 'Продукт', 'Приготовление', 'Перепутали', 'Сервис', 'Опоздание', 'Всего']
    else:
        headers = ['Ресторан', 'Вкус', 'Быстрая доставка', 'Вежливый персонал', 'Качество', 'Атмосфера', 'Всего']
    
    # Ищем строки с данными
    data_rows = []
    for idx in range(start_row + 2, len(df)):  # +2 чтобы пропустить заголовок
        row = df.iloc[idx]
        restaurant = str(row[0]).strip()
        
        # Пропускаем пустые строки и итоги
        if not restaurant or 'Всего' in restaurant or 'Анализ' in restaurant:
            continue
        
        # Проверяем, что это название ресторана
        if restaurant and len(restaurant) > 2:
            try:
                data_rows.append({
                    'Ресторан': restaurant,
                    headers[1]: row[1] if pd.notna(row[1]) else 0,
                    headers[2]: row[2] if pd.notna(row[2]) else 0,
                    headers[3]: row[3] if pd.notna(row[3]) else 0,
                    headers[4]: row[4] if pd.notna(row[4]) else 0,
                    headers[5]: row[5] if pd.notna(row[5]) else 0,
                    headers[6]: row[6] if pd.notna(row[6]) else 0,
                })
            except:
                continue
    
    return pd.DataFrame(data_rows)

# ==========================================
# 🍕 ГЛАВНАЯ ФУНКЦИЯ ГЕНЕРАЦИИ
# ==========================================
def generate_flexible_presentation(
    all_dataframes: Dict[str, Dict[str, pd.DataFrame]], 
    slides_structure: List[str],
    period_text: str,
    theme: str = "Классическая (Томат + Базилик)",
    output_path: str = "Гибкая_презентация.pptx"
):
    """
    Генерирует презентацию с умным парсингом данных КР.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Выбираем цвета темы
    if "Оливковая" in theme:
        primary_color = PJColors.OLIVE_GREEN
        secondary_color = PJColors.BASIL_GREEN
    elif "Сырная" in theme:
        primary_color = PJColors.CHEESE_YELLOW
        secondary_color = PJColors.TOMATO_RED
    else:
        primary_color = PJColors.TOMATO_RED
        secondary_color = PJColors.BASIL_GREEN
    
    # --- СЛАЙД 1: ТИТУЛЬНЫЙ ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, primary_color)
    
    _add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
              "КЛИЕНТСКИЙ РЕЙТИНГ", PJFonts.HEADLINE, 64, PJColors.WHITE, bold=True, align=PP_ALIGN.CENTER)
    
    _add_text(slide, Inches(1), Inches(3.2), Inches(11), Inches(1),
              period_text, PJFonts.BODY, 36, PJColors.CHEESE_YELLOW, align=PP_ALIGN.CENTER)
    
    _add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(1),
              PJSlogans.EMOTION, PJFonts.BODY, 28, PJColors.WHITE, align=PP_ALIGN.CENTER)
    
    # --- ГЕНЕРАЦИЯ СЛАЙДОВ ПО СТРУКТУРЕ ---
    file_list = list(all_dataframes.keys())
    
    for slide_idx, slide_title in enumerate(slides_structure):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_slide_bg(slide, PJColors.WHITE)
        
        # Заголовок слайда
        _add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                  slide_title, PJFonts.HEADLINE, 40, secondary_color, bold=True)
        
        # Определяем тип слайда по ключевым словам
        slide_lower = slide_title.lower()
        
        # Получаем данные из первого файла
        if file_list:
            first_file = file_list[0]
            parsed_data = all_dataframes[first_file]
            
            # Определяем, какие данные показать
            df_to_show = None
            
            if 'общ' in slide_lower or 'показател' in slide_lower or 'сеть' in slide_lower:
                # Сводный слайд - показываем средние рейтинги
                df_to_show = _create_summary_table(parsed_data)
            
            elif 'спб' in slide_lower and 'оценк' in slide_lower:
                # Оценки СПб - берем все источники
                df_to_show = _combine_city_ratings(parsed_data, 'спб')
            
            elif 'тюмень' in slide_lower and 'оценк' in slide_lower:
                # Оценки Тюмень
                df_to_show = _combine_city_ratings(parsed_data, 'тюмень')
            
            elif 'жалоб' in slide_lower:
                # Жалобы
                if 'жалобы' in parsed_data:
                    df_to_show = parsed_data['жалобы']
            
            elif 'позитив' in slide_lower or 'положительн' in slide_lower:
                # Позитивные отзывы
                if 'позитив' in parsed_data:
                    df_to_show = parsed_data['позитив']
            
            else:
                # По умолчанию - первый доступный блок
                for key in ['сайт_спб', 'агрегаторы_спб', 'жалобы', 'позитив']:
                    if key in parsed_data and not parsed_data[key].empty:
                        df_to_show = parsed_data[key]
                        break
            
            # Отображаем таблицу
            if df_to_show is not None and not df_to_show.empty:
                _add_data_table(slide, df_to_show, secondary_color)
            else:
                _add_text(slide, Inches(1), Inches(2), Inches(11), Inches(3),
                         "Нет данных для отображения",
                         PJFonts.BODY, 24, PJColors.OLIVE_GREEN, align=PP_ALIGN.CENTER)
        
        # Номер слайда
        _add_text(slide, Inches(12), Inches(7), Inches(1), Inches(0.5),
                  f"{slide_idx + 2}/{len(slides_structure) + 1}", 
                  PJFonts.BODY, 12, PJColors.OLIVE_GREEN)
    
    # --- ФИНАЛЬНЫЙ СЛАЙД ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, secondary_color)
    
    _add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.5),
              PJSlogans.MAIN, PJFonts.HEADLINE, 54, PJColors.WHITE, bold=True, align=PP_ALIGN.CENTER)
    
    _add_text(slide, Inches(1), Inches(4.5), Inches(11), Inches(1),
              PJSlogans.CARE, PJFonts.BODY, 28, PJColors.CHEESE_YELLOW, align=PP_ALIGN.CENTER)
    
    prs.save(output_path)
    return output_path

def _create_summary_table(parsed_data: Dict) -> pd.DataFrame:
    """Создает сводную таблицу со средними рейтингами."""
    summary = []
    
    for key in ['сайт_спб', 'сайт_тюмень', 'агрегаторы_спб', 'агрегаторы_тюмень', 
                'геосервисы_спб', 'геосервисы_тюмень']:
        if key in parsed_data and not parsed_data[key].empty:
            df = parsed_data[key]
            # Берем строку "Итого" если есть, иначе считаем сами
            if 'Итого' in str(df['Ресторан'].values):
                total_row = df[df['Ресторан'].str.contains('Итого', na=False)]
                if not total_row.empty:
                    avg = total_row.iloc[0]['Средний рейтинг']
                else:
                    avg = df['Средний рейтинг'].mean()
            else:
                avg = df['Средний рейтинг'].mean()
            
            source = key.split('_')[0]
            city = key.split('_')[1].upper()
            summary.append({'Источник': source, 'Город': city, 'Средний рейтинг': round(avg, 2)})
    
    return pd.DataFrame(summary)

def _combine_city_ratings(parsed_data: Dict, city: str) -> pd.DataFrame:
    """Объединяет оценки по городу из всех источников."""
    combined = pd.DataFrame()
    
    for source in ['сайт', 'агрегаторы', 'геосервисы']:
        key = f'{source}_{city}'
        if key in parsed_data and not parsed_data[key].empty:
            df = parsed_data[key].copy()
            df['Источник'] = source.capitalize()
            if combined.empty:
                combined = df
            else:
                combined = pd.concat([combined, df], ignore_index=True)
    
    return combined

def _add_data_table(slide, df: pd.DataFrame, accent_color):
    """Добавляет таблицу с данными на слайд."""
    rows = min(len(df) + 1, 15)  # Максимум 14 строк данных + заголовок
    cols = len(df.columns)
    
    if cols == 0 or rows <= 1:
        return
    
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), 
                                        Inches(12), Inches(5))
    table = table_shape.table
    
    # Заголовки
    for col_idx, col_name in enumerate(df.columns):
        _format_cell(table.cell(0, col_idx), str(col_name), 
                   PJFonts.HEADLINE, 16, PJColors.WHITE, accent_color, bold=True)
    
    # Данные
    for row_idx, (_, row) in enumerate(df.head(14).iterrows()):
        bg = PJColors.DOUGH_BEIGE if row_idx % 2 == 0 else PJColors.WHITE
        for col_idx, value in enumerate(row):
            if col_idx < cols:
                _format_cell(table.cell(row_idx + 1, col_idx), str(value),
                           PJFonts.BODY, 14, PJColors.BLACK, bg)
